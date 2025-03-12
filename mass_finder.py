import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
from scipy.signal import find_peaks
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

import all_elements
import mass_calculator
import spectrum_simulator


# Variables:
# root is the path where the data folder will be located
root = r'X:\My Documents\HRMS\data'

# name is the name of the folder where the txt files will be located
name="testbatch_HHA" 

# allowed_deviation is how much the expected mass may differ from the found mass in ppm
allowed_deviation=4

# powerpoint_information is some text at the bottom of the powerpoint file.
powerpoint_information = "Measured on the Thermo Scientific Orbitrap Exploris 480"

sequence_data_path = r"D:\Data"
sequence_method_path = r"D:\01-Instrument-Methods"

def extract_data(data_text, charge):
    '''
    Looks in the data text for the spectrum with the highest total ion count of the given charge. Extracts the spectrum as x and y coordinates (masses and signals).
    Also calculates nomalized values for the signals, the highest value being set to 100.
    
    Args:
        data_text (list[str]): All lines extracted from mass file converted to txt.
        charge (str): Either + or -, depending on what ionization mode to check.
        
    Returns:
        masses (np.array[float]): Array of all masses from the spectrum.
        signals (np.array[float]): Array of all signals from the spectrum.
        signals_normalized (np.array[float]): Array of all normalized signals.
    
    '''
    
    scans = []
    for count, line in enumerate(data_text):
        if "index: " in line:
            scans.append(count)
            if len(scans) > 1:
                break
    scan_size = scans[1] - scans[0]
    all_scans = {}
    if charge == "+":
        select_scan = "positive scan"
    elif charge == "-":
        select_scan = "negative scan"
    for count, line in enumerate(data_text):
        if "index: " in line:
            right_mode = False
            total_ion_current = None
            data = []
            for line2 in data_text[count:count+scan_size]:
                if select_scan in line2:
                    right_mode = True
                elif "total ion current," in line2:
                    total_ion_current = float(line2.split()[-1])
                elif "binary: " in line2:
                    data.append(line2)
            if right_mode and total_ion_current:
                all_scans[total_ion_current] = data

    masses_and_signals = all_scans[max(all_scans.keys())]
    masses = masses_and_signals[0].split()[2:-1]
    signals = masses_and_signals[1].split()[2:-1]

    for i, (mass, signal) in enumerate(zip(masses, signals)):
        masses[i] = float(mass)
        signals[i] = float(signal)

    masses = np.array(masses)
    signals = np.array(signals)
    signalmax=signals.max()
    signals_normalized = (signals/signalmax)*100
    
    return masses,signals, signals_normalized
    

class Sample:
    def __init__(self, sample_name, molecular_formula, e_mail):
        self.sample_name=sample_name
        self.molecular_formula=molecular_formula
        self.theoretical_mass=mass_calculator.mass_and_element_calculator(molecular_formula, 0)[2]
        self.found_mass=0
        self.ion = 0
        self.ion_text = ""
        self.relative_abundance = 0
        self.success = False
        self.e_mail = e_mail


def plot(sample, masses, signals, signals_normalized, figure_name1, figure_name2, charge):
    '''
    Plots the spectrum with peak annotation and plots a zoom in spectrm on the area of the expected mass. Saves the two pictures at the specified paths.

    Args:
        sample (Sample): The sample containing the information to be plotted.
        masses (np.array[float]): Array of all masses from the spectrum.
        signals (np.array[float]): Array of all signals from the spectrum.
        signals_normalized (np.array[float]): Array of all normalized signals.
        figure_name1 (str): The name of the figure containing the whole spectrum.
        figure_name2 (str): The name of the figure containing the zoom in spectrum.
        charge: Either + or -, depending on what ionization was used.
    
    Returns:
        None
    '''
    peak, height=spectrum_simulator.spectrum_simulator(sample.ion, charge)
    heightnorm=[]
    heightmax=max(height)
    
    for i in height:
        sig=(i/heightmax)*100
        heightnorm.append(sig)
        
    mass=mass_calculator.mass_and_element_calculator(sample.ion, charge)[2]
    low=mass
    for i in range(len(heightnorm)):
        if heightnorm[i]>2:
            if peak[i] < low:
                low=peak[i]
    low-=0.3  
    up=0
    for i in range(len(heightnorm)):
        if heightnorm[i]>2:
            if peak[i] > up:
                up=peak[i]
    up+=0.3

    zoom_masses = [0]
    zoom_signals = [0]
    zoom_signals_norm = []
    
    
    for i in range(len(masses)):
        if masses[i] > low and masses[i] < up:
            zoom_masses.append(masses[i])
            zoom_signals.append(signals[i])
            
    zoomSignalMax = max(zoom_signals)
    if zoomSignalMax == 0:
        zoomSignalMax = 1
    for i in zoom_signals:
        zoom_signals_norm.append((i/zoomSignalMax)*100)
    zoom_masses.append(10000)
    zoom_signals_norm.append(0)
    
    fig, axs1 = plt.subplots(2, sharex=True, sharey=True, figsize=(12,6), gridspec_kw={'hspace': 0})    
    plt.xlim(masses[0],masses[-1])
    plt.ylim(0,110)
    
    peaks,_= find_peaks(signals_normalized, height=10)
    zoom_peaks,_ = find_peaks(zoom_signals_norm, height=3)
   
    fig2, axs2 = plt.subplots(2, sharex=True, sharey=True, figsize=(12,6), gridspec_kw={'hspace': 0})
    plt.xlim(low,up)
    plt.ylim(0,110) 
    
    all_axs = np.concatenate((axs1, axs2))
    
    all_axs[0].plot(masses,signals_normalized, color='k', linewidth = 0.3)
    markerline, stemlines, baseline = all_axs[1].stem(peak,heightnorm,markerfmt=" ", basefmt=" ", linefmt = 'k')
    plt.setp(stemlines, linewidth=0.3)

    
    all_axs[2].plot(zoom_masses,zoom_signals_norm, color='k', linewidth = 0.3)
    markerline, stemlines, baseline = all_axs[3].stem(peak,heightnorm,markerfmt=" ", basefmt=" ", linefmt = 'k')
    plt.setp(stemlines, linewidth=0.3)   
    
    for i in all_axs:
        i.set_xlabel("m/z")
        i.set_ylabel("Relative Abundance")
        i.spines['top'].set_visible(False)
        i.spines['right'].set_visible(False)
    
    

    all_axs[0].text(0.85, 0.9, 'Measured', fontsize=12, transform=all_axs[0].transAxes)
    all_axs[0].text(0.85, 0.8, sample.sample_name, fontsize=12, transform=all_axs[0].transAxes)
    all_axs[1].text(0.85, 0.9, 'Theoretical', fontsize=12, transform=all_axs[1].transAxes)
    all_axs[1].text(0.85, 0.8, sample.ion + charge, fontsize=12, transform=all_axs[1].transAxes)
    
    all_axs[2].text(0.85, 0.9, 'Measured', fontsize=12, transform=all_axs[2].transAxes)
    all_axs[2].text(0.85, 0.8, sample.sample_name, fontsize=12, transform=all_axs[2].transAxes)
    all_axs[3].text(0.85, 0.9, 'Theoretical', fontsize=12, transform=all_axs[3].transAxes)  
    all_axs[3].text(0.85, 0.8, sample.ion + charge, fontsize=12, transform=all_axs[3].transAxes)


    for i in peaks:
        all_axs[0].plot(masses[i], signals_normalized[i])
        text= "{:.5f}".format(masses[i])
        all_axs[0].annotate(text, xy=(masses[i]-20, signals_normalized[i]+2), fontsize = 7)
    
        
    for i in range(len(heightnorm)):
        if heightnorm[i] > 50:
            all_axs[1].plot(peak[i], heightnorm[i])
            text= "{:.5f}".format(peak[i])
            all_axs[1].annotate(text, xy=(peak[i]-20, heightnorm[i]+2), fontsize = 7)
    
    for i in zoom_peaks:
        if low < zoom_masses[i] < up:
            all_axs[2].plot(zoom_masses[i], zoom_signals_norm[i])
            text= "{:.5f}".format(zoom_masses[i])
            all_axs[2].annotate(text, xy=(zoom_masses[i], zoom_signals_norm[i]), fontsize = 7)
        
    for i in range(len(heightnorm)):
        if heightnorm[i] > 3:
            all_axs[3].plot(peak[i], heightnorm[i])
            text= "{:.5f}".format(peak[i])
            all_axs[3].annotate(text, xy=(peak[i], heightnorm[i]), fontsize =7)
        
    fig.savefig(figure_name1, dpi=199)
    fig2.savefig(figure_name2, dpi=199)

def ion_calculator(adduct, molecular_formula, allowed_deviation, charge):
    '''
    Calculates the mass of an ion given a molecular formula and an adduct. Also calculated upper and lower limit of the range within the measured mass should fall.

    Args:
        adduct (str): One or more atoms that are being added to the molecular formula. Can also be "" if no adduct is present or "-1" if one H should be subtracted.
        molecular_formula (str): The molecular formula to get the ion mass from.
        allowed_deviation (int): The allowed deviation of the found mass in ppm.
        charge (str): Either + or - depending on what ionization mode was used.
    
    Returns:
        calculated_mass (float): The calculated mass of the ion.
        lower (float): The lower mass limit.
        upper (float): The upper mass limit.
    '''
    if adduct == "-1":
        molecular_formula = mass_calculator.minus_H(molecular_formula)
    molecular_mass = mass_calculator.mass_and_element_calculator(molecular_formula, charge)[2]
    if adduct in all_elements.atoms:
        calculated_mass = molecular_mass+all_elements.atoms[adduct]
    else:
        calculated_mass=molecular_mass
    deviation = (calculated_mass/1000000)*allowed_deviation
    lower = calculated_mass-deviation
    upper = calculated_mass+deviation
    return calculated_mass, lower, upper

def match(sample, allowed_deviation, masses, signals, signals_normalized, charge):
    '''
    Checks the spectrum for the expected mass with multiple adducts. Will update the Sample instance with the information.

    Args:
        sample (Sample): The sample containing the information to be plotted.
        allowed_deviation (int): The allowed deviation of the found mass in ppm.
        masses (np.array[float]): Array of all masses from the spectrum.
        signals (np.array[float]): Array of all signals from the spectrum.
        signals_normalized (np.array[float]): Array of all normalized signals.
        charge (str): Either + or - depending on what ionization mode was used.

    Returns:
        succes (bool): Whether the expected mass was found in the spectrum.

    '''
    
    if charge == "+":
        possible_adducts = ["H", "Na", "NH4", ""]
    elif charge == "-":
        possible_adducts = ["", "-1"]
        
    succes=0
    for adduct in possible_adducts:
            calculated_mass, lower,upper=ion_calculator(adduct, sample.molecular_formula, allowed_deviation, charge)
            for j in range(len(masses)):
                if lower < masses[j] < upper:
                    if signals_normalized[j] > 1:
                        while True:
                            if signals_normalized[j+1] > signals_normalized[j]:
                                j += 1
                            else:
                                break
                        if signals_normalized[j] > sample.relative_abundance:
                            sample.found_mass = masses[j]
                            sample.relative_abundance = signals_normalized[j]
                            if adduct == "-1":
                                sample.ion = mass_calculator.minus_H(sample.molecular_formula)
                            else:
                                sample.ion = sample.molecular_formula+adduct
                            if charge == "+":
                                if adduct == "":
                                    sample.ion_text = "[M]+"
                                else:
                                    sample.ion_text = "[M+{}]+".format(adduct)
                            else:
                                if adduct == "-1":
                                    sample.ion_text = "[M-H]-"
                                elif adduct == "":
                                    sample.ion_text = "[M]-"
                            sample.theoretical_mass = mass_calculator.mass_and_element_calculator(sample.ion, charge)[2]
                            print(sample.sample_name + "(" + charge + ")" +":  " + "Succes with {}".format(adduct) + "  (" + str(calculated_mass) + ")" + "  (" + str(masses[j]) + ")" + "  (" + str(signals[j]) + ")")
                            succes=1
                            sample.success = True
        
    if succes==0:
        print(sample.sample_name + "(" + charge + ")" + ":  " + "Fail")
        if charge == "+":
            sample.ion = sample.molecular_formula + "H"
        if charge == "-":
            sample.ion = mass_calculator.minus_H(sample.molecular_formula)
            
    return succes

def create_powerpoint(sample, powerpoint_path, figure_name1, figure_name2, charge, succes):
        '''
        Creates a powerpoint file containing the spectrum figures and some additional information.

        Args:
            sample (Sample): The sample containing the information to be plotted.
            powerpoint_path (str): The path where the powerpoint file will be saved.
            figure_name1 (str): The name of the figure containing the whole spectrum.
            figure_name2 (str): The name of the figure containing the zoom in spectrum.
            charge (str): Either + or - depending on what ionization mode was used.
            success (bool): Whether or not the expected mass was found in the spectrum.

        Returns:
            None
        '''
        powerpoint_report=Presentation()
        difference = round(abs(((sample.found_mass - sample.theoretical_mass)/sample.theoretical_mass)*1000000), 2)
        
        Layout = powerpoint_report.slide_layouts[1] 
        first_slide = powerpoint_report.slides.add_slide(Layout)
        title = first_slide.shapes.title.text_frame.paragraphs[0]
        title.alignment = PP_ALIGN.LEFT
        title.font.size = Cm(1)
        title.text = sample.sample_name
        
        figure1 = first_slide.shapes.add_picture(figure_name1, Cm(1), Cm(4), width=Cm(24))
        location_information_frame = first_slide.shapes.add_textbox(Cm(1), Cm(17), Cm(17), Cm(2))
        location_information = location_information_frame.text_frame.paragraphs[0]
        location_information.font.size = Cm(0.5)
        location_information.text = powerpoint_information
        
        sample_information_frame = first_slide.shapes.add_textbox(Cm(15), Cm(1), Cm(5), Cm(2))
        sample_information = sample_information_frame.text_frame.paragraphs[0]
        sample_information.font.size = Cm(0.5)
        if charge == "+":
            sample_information.text = f"Ionization mode: ESI positive \nFound mass: {round(sample.found_mass, 4):.4f}\nRelative abundance: {sample.relative_abundance:.1f}%\nTheoretical mass: {round(sample.theoretical_mass, 4):.4f}\nDifference (ppm): {difference:.2f}"
        if charge == "-":
            sample_information.text = f"Ionization mode: ESI negative \nFound mass: {round(sample.found_mass, 4):.4f}\nRelative abundance: {sample.relative_abundance:.1f}%\nTheoretical mass: {round(sample.theoretical_mass, 4):.4f}\nDifference (ppm): {difference:.2f}"
        second_slide = powerpoint_report.slides.add_slide(Layout)
        title2 = second_slide.shapes.title.text_frame.paragraphs[0]
        title2.alignment = PP_ALIGN.LEFT
        title2.font.size = Cm(1)
        title2.text = sample.sample_name + " (zoomed in)"
        figure2 = second_slide.shapes.add_picture(figure_name2, Cm(1), Cm(4), width=Cm(24))
        
        person_path = os.path.join(powerpoint_path, sample.e_mail)
        failed_path = powerpoint_path
        results_file = os.path.join(person_path, "Found masses.csv")
        txt_file = os.path.join(person_path, "Found masses.txt")
        
        if not os.path.exists(person_path):
            os.mkdir(person_path)
            file = open(results_file, "w")
            file.write("Mode,Sample name,Formula,Ion,Expected mass,Found mass,Deviation (ppm)\n")
            file.close()
            file = open(txt_file, "w")
            file.close()
        if succes == 1:
            powerpoint_report.save(os.path.join(person_path, f"{sample.sample_name} ({charge}).pptx"))
            csv_text = ["ESI", sample.sample_name, sample.ion, sample.ion_text, str(round(sample.theoretical_mass, 4)), str(round(sample.found_mass, 4)), str(difference)]
            txt_text = f"{sample.sample_name}: HRMS (ESI) m/z for {sample.ion_text} calculated for {sample.ion} {str(round(sample.theoretical_mass, 4))}; found {str(round(sample.found_mass, 4))}.\n"
            file = open(results_file, "a")
            file.write(",".join(csv_text) + "\n")
            file.close()
            file = open(txt_file, "a")
            file.write(txt_text)
            file.close()
        else:
            powerpoint_report.save(os.path.join(failed_path, sample.sample_name + f"(failed) {charge}).pptx"))
        
        

def sequence_maker(dataframe, data_folder, name):
    '''
    Creates a new sequence for the failed samples to be measured in APCI mode.

    Args:
        dataframe (pd.DataFrame): The old sequence to be converted to APCI sequence.
        data_folder (str): The path of the folder where the sequence will be saved.
        name (str): The name of the main datafolder where the txt files are located.

    Returns:
        None
    '''
    
    out=dataframe

    out["Path"]=os.path.join(sequence_data_path, f"{str(name)}APCIpos")
    
    for i in range(len(out["File Name"])):
        if dataframe.iloc[i,2] == os.path.join(sequence_method_path, "ESI posneg solvC 5min 100-1000 noUV"):
            out.iloc[i,2]= os.path.join(sequence_method_path, "APCI pos solvC 5min 100-1000 noUV")
        elif dataframe.iloc[i,2] == os.path.join(sequence_method_path, "ESI posneg solvC 5min 250-2500 noUV"):
            out.iloc[i,2]= os.path.join(sequence_method_path, "APCI pos solvC 5min 250-2500 noUV")
        else:
            out.iloc[i,2]= r""
    
      
    file_name = os.path.join(data_folder, "sequence_APCI.csv")
    seq=open(file_name, "w+")          
    seq.write("Bracket Type=4\n")
    seq.close()
    out.to_csv(file_name, mode='a', index=False)
    


def main():
    '''
    Reads in a file called samples.csv where information about samples is located. Matches these sample names to file names with .txt extension which are .raw files converted by the MSconvert software by ProteoWizard.
    Looks for the expected masses in the spectra and creates reports as .png and .pptx files. Creates a new APCI sequence for the failed samples if a sequence.csv file is provided.

    Args:
        None

    Returns:
        None
    '''
    data_folder = os.path.join(root, str(name))
    start_time = datetime.now()
    csv_file = os.path.join(data_folder, "samples.csv")
    figures_path = os.path.join(data_folder, "figures")
    if not os.path.exists(figures_path):
        os.mkdir(figures_path)
    powerpoint_path = os.path.join(data_folder, "powerpoints")
    if not os.path.exists(powerpoint_path):
        os.mkdir(powerpoint_path)
    df = pd.read_csv(csv_file, header=None)
    sequence_path = os.path.join(data_folder, "sequence.csv")
    if os.path.exists(sequence_path):
        sequence = pd.read_csv(os.path.join(data_folder, "sequence.csv"), skiprows=1)
        sequence["Success"] = ""
    else:
        sequence = pd.DataFrame()
    
    for i in range(len(df[8])):
        df.at[i,8]=str(df[8][i]).replace(" ","_")
        df.at[i,8]=str(df[8][i]).replace("*","_")
        df.at[i,8]=str(df[8][i]).replace(".","_")
        df.at[i,6]=str(df[6][i]).replace("+","")
        df.at[i,6]=str(df[6][i]).replace("-","")
        df.at[i,6]=str(df[6][i]).replace(" ","")
        
    failed_samples = df.copy()
    failed_samples["Succes"] = ""
        
    for i in range(len(df[0])):
        active_sample = Sample(str(df.iloc[i,8]), molecular_formula=str(df.iloc[i,6]), e_mail=str(df.iloc[i, 4]))
        try:
            file = open((os.path.join(data_folder, active_sample.sample_name + ".txt.")), 'r')
        except:
            print("file {} failed to open".format(active_sample.sample_name))
            continue
        data_text = file.readlines()
        file.close()
        charges = ["+", "-"]
        for charge in charges:
            masses_pos, signals_pos, signals_normalized_pos = extract_data(data_text, charge)
            succes = match(active_sample, allowed_deviation, masses_pos, signals_pos, signals_normalized_pos, charge)
            figure_name1 = os.path.join(data_folder, "figures", active_sample.sample_name + ".png")
            figure_name2 = os.path.join(data_folder, "figures", active_sample.sample_name + "_zoom.png")
            plot(active_sample, masses_pos, signals_pos, signals_normalized_pos, figure_name1, figure_name2, charge)
            create_powerpoint(active_sample, powerpoint_path, figure_name1, figure_name2, charge, succes)
        if active_sample.success == False:
            failed_samples.loc[i, "Succes"] = "Failed"
        else:
            if not sequence.empty:
                for j in range(len(sequence["File Name"])):
                    if active_sample.sample_name == sequence["File Name"][j]:
                        sequence.at[j, "Success"] = "Yes"
    if not sequence.empty:
        sequence = sequence[sequence["Success"] != "Yes"]     
        sequence = sequence.drop("Success", axis=1)
    failed_samples = failed_samples[failed_samples["Succes"] == "Failed"]
    APCI_folder = os.path.join(data_folder, "APCI")
    if not os.path.exists(APCI_folder):
        os.mkdir(APCI_folder)
    file_name = os.path.join(APCI_folder, "APCI samples.csv")
    failed_samples_csv = open(file_name, "w+")
    failed_samples_csv.close()
    failed_samples.to_csv(file_name, index=False, header=False)
    
    if not sequence.empty:
        sequence_maker(sequence, APCI_folder, name)
    
    print("Total time: ")
    print(datetime.now()-start_time)


if __name__ == "__main__":
    main()

